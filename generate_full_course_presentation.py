from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Couleurs Kadea Academy
KADEA_RED = RGBColor(237, 28, 36)
ANTHRACITE = RGBColor(40, 40, 40)
WHITE = RGBColor(255, 255, 255)

def set_shape_text(shape, text, color=ANTHRACITE, font_size=Pt(16), bold=False, alignment=PP_ALIGN.LEFT):
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.size = font_size
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment

def add_slide_with_content(prs, title_text, text_blocks, viz_text="", speaker_notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, title_text, KADEA_RED, Pt(26), True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for block in text_blocks:
        p = tf.add_paragraph()
        p.text = block['title']
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = ANTHRACITE if not block.get('highlight') else KADEA_RED

        for bullet in block['bullets']:
            p_bullet = tf.add_paragraph()
            p_bullet.text = bullet
            p_bullet.font.size = Pt(14)
            p_bullet.level = 1

    # Visual Placeholder
    if viz_text:
        viz_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(1.5), Inches(3.5), Inches(5))
        viz_box.fill.solid()
        viz_box.fill.fore_color.rgb = WHITE
        viz_box.line.color.rgb = KADEA_RED
        viz_box.line.width = Pt(2)
        set_shape_text(viz_box, viz_text, ANTHRACITE, Pt(12), False, PP_ALIGN.CENTER)

    # Add Speaker Notes
    if speaker_notes and slide.notes_slide:
        text_frame = slide.notes_slide.notes_text_frame
        text_frame.text = speaker_notes

def create_full_presentation():
    prs = Presentation()

    # Slide 1: Titre
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_shape_text(slide.shapes.title, "Ingénierie de la Donnée : Modélisation Avancée et SQL", KADEA_RED, Pt(32), True, PP_ALIGN.CENTER)
    set_shape_text(slide.placeholders[1], "Architecture Relationnelle, Optimisation et Analyse de Données", ANTHRACITE, Pt(24), False, PP_ALIGN.CENTER)
    if slide.notes_slide:
        slide.notes_slide.notes_text_frame.text = "Bienvenue dans ce module avancé. L'objectif n'est pas seulement d'écrire du SQL, mais de concevoir des architectures robustes et performantes. Nous allons explorer les fondamentaux théoriques, puis plonger dans la modélisation UML, l'optimisation des performances, et les requêtes analytiques avancées."

    # Slide 2: Généralités
    add_slide_with_content(
        prs,
        "1. Les Fondations de l'Ingénierie de Données",
        [
            {"title": "Au-delà du Tableur", "bullets": ["Les bases de données séparent le stockage physique de la logique applicative.", "Objectif : Indépendance physique et logique des données."]},
            {"title": "Le Rôle du SGBDR (RDBMS)", "bullets": ["Le SGBDR n'est pas qu'un espace de stockage, c'est un moteur de règles.", "Gestion stricte de la concurrence (verrous, isolation).", "Gestion de la reprise sur panne (Write-Ahead Logging)."]},
        ],
        "[VISUEL KADEA]\nArchitecture Client-Serveur: Couche Applicative -> Moteur SGBD (Optimiseur de requêtes) -> Stockage Disque.",
        speaker_notes="Le SGBD est souvent réduit à tort à un gros disque dur. Insistez sur le fait que le SGBD est le garant de l'intégrité métier. Sans lui, des transactions concurrentes corrompraient la base instantanément. Mentionnez le 'Write-Ahead Logging' (WAL) qui permet aux banques de ne jamais perdre une transaction même en cas de coupure de courant."
    )

    # Slide 3: Historique
    add_slide_with_content(
        prs,
        "2. Les Différents Types de SGBD (Historique)",
        [
            {"title": "Années 1960", "bullets": ["Hiérarchiques : données en arborescence.", "Réseaux : arborescence avec raccourcis."]},
            {"title": "Années 1970 - Modèle Relationnel", "bullets": ["1970 : Edgar F. Codd invente le modèle relationnel.", "Données organisées en tableaux liés entre eux. Langage SQL standardisé en 1974."]},
            {"title": "Années 2000 - NoSQL", "highlight": True, "bullets": ["Not Only SQL : Bases Clé/Valeur, Orientées Documents, Graphes.", "Pour le Big Data et la très haute disponibilité."]}
        ],
        "[VISUEL KADEA]\nFrise chronologique des SGBD, mettant en avant l'apparition du Modèle Relationnel et du NoSQL."
    )

    # Slide 4: Transactions ACID
    add_slide_with_content(
        prs,
        "3. L'Architecture Transactionnelle : ACID",
        [
            {"title": "La Transaction de bout en bout", "bullets": ["L'unité logique de traitement garantissant l'intégrité systémique."]},
            {"title": "Les Principes Fondamentaux", "highlight": True, "bullets": [
                "Atomicité (All-or-Nothing) : Exécution complète ou Rollback total.",
                "Cohérence : Application stricte des contraintes (Constraints & Triggers).",
                "Isolation : Prévention des lectures sales (Dirty Reads) via les niveaux d'isolation.",
                "Durabilité : Persistance garantie même en cas d'arrêt brutal (Crash)."
            ]}
        ],
        "[VISUEL KADEA]\nIcône d'un coffre-fort avec les 4 piliers de l'ACID. Indiquer 'Commit' et 'Rollback'.",
        speaker_notes="Le concept d'isolation est le plus subtil. Expliquez les 'Niveaux d'Isolation' (Read Uncommitted, Read Committed, Repeatable Read, Serializable). Un DBA doit souvent arbitrer entre le plus haut niveau d'isolation (sécurité maximale mais performances réduites à cause des verrous) et des niveaux inférieurs plus rapides."
    )

    # Slide 5: Le Modèle Relationnel (Concepts)
    add_slide_with_content(
        prs,
        "4. Le Modèle Relationnel : Fondements",
        [
            {"title": "Relations (Tables)", "bullets": ["L'information est stockée dans des tables constituées de colonnes (attributs) et de lignes (enregistrements).", "Chaque colonne possède un domaine de valeurs précis (type, longueur, contraintes)."]},
            {"title": "Le Système de Clés", "highlight": True, "bullets": [
                "Clé Primaire : Identifiant unique d'une ligne dans une table.",
                "Clé Étrangère : Colonne pointant vers la clé primaire d'une autre table, créant le lien."
            ]}
        ],
        "[VISUEL KADEA]\nIllustration de deux tables (Ville et Personne). Un rayon laser relie la Clé Primaire 'Code_Ville' à la Clé Étrangère de la table Personne."
    )

    # Slide 6: Le Modèle Relationnel (Cardinalités)
    add_slide_with_content(
        prs,
        "5. Associations et Cardinalités",
        [
            {"title": "Types de Liens", "bullets": [
                "1-1 : Un pays a une seule capitale.",
                "1-N : Un pays possède plusieurs villes, une ville appartient à un seul pays."
            ]},
            {"title": "Le Cas Complexe : Plusieurs-à-Plusieurs (N-M)", "highlight": True, "bullets": [
                "Ex: Une personne travaille dans plusieurs entreprises, une entreprise emploie plusieurs personnes.",
                "Solution : Créer une 'Table d'Association' (ex: Travail) contenant les clés étrangères des deux tables."
            ]}
        ],
        "[VISUEL KADEA]\nSchéma expliquant la résolution d'une relation N-M par une table de jonction.",
        speaker_notes="Insistez sur la Table de Jonction (Junction Table). Une erreur classique de débutant est d'essayer de stocker des valeurs multiples séparées par des virgules dans une seule colonne. Cela casse complètement le modèle relationnel et les performances de requêtage."
    )

    # NOUVELLE SLIDE : NORMALISATION
    add_slide_with_content(
        prs,
        "6. Théorie de l'Architecture : La Normalisation",
        [
            {"title": "Pourquoi normaliser ?", "bullets": [
                "Éviter les anomalies d'insertion, de mise à jour et de suppression.",
                "Éliminer la redondance des données pour garantir une source unique de vérité."
            ]},
            {"title": "Les Formes Normales (Normal Forms)", "highlight": True, "bullets": [
                "1NF : Valeurs atomiques (pas de listes dans une colonne).",
                "2NF : 1NF + Tout attribut non clé dépend de la totalité de la clé primaire.",
                "3NF : 2NF + Pas de dépendance transitive (un attribut non clé ne doit pas dépendre d'un autre attribut non clé).",
                "BCNF (Boyce-Codd) : Version stricte de la 3NF."
            ]}
        ],
        "[VISUEL KADEA]\nTableau montrant une donnée 'dénormalisée' se divisant proprement en trois tables normalisées (1NF -> 3NF).",
        speaker_notes="La 3ème Forme Normale (3NF) est le standard d'or de l'industrie. La phrase mnémotechnique de Bill Kent : 'Every non-key attribute must provide a fact about the key, the whole key, and nothing but the key, so help me Codd.' Mentionnez que dans les entrepôts de données (Data Warehouses) pour la BI, on effectue souvent le processus inverse : la 'Dénormalisation', pour accélérer la lecture."
    )

    # NOUVELLE SLIDE : INDEX ET PERFORMANCES
    add_slide_with_content(
        prs,
        "7. Performance : Indexation et Plans d'Exécution",
        [
            {"title": "Les Mécanismes d'Indexation", "bullets": [
                "Un index accélère la lecture mais ralentit l'écriture (INSERT/UPDATE).",
                "Arbres B-Tree (Balanced Tree) : Standard pour les recherches de plages.",
                "Index Hash : Optimisés pour les correspondances exactes."
            ]},
            {"title": "Query Execution Plan (EXPLAIN)", "highlight": True, "bullets": [
                "L'optimiseur (Query Planner) décide comment lire les données.",
                "Full Table Scan : Lecture séquentielle (très coûteuse).",
                "Index Seek / Scan : Accès direct via l'index O(log N)."
            ]}
        ],
        "[VISUEL KADEA]\nArbre B-Tree inversé cherchant la valeur 42, contre un tableau scanné ligne par ligne.",
        speaker_notes="Si une requête SQL est lente, la solution n'est presque jamais 'le SGBD est mauvais'. C'est généralement un manque d'index ou une requête mal formulée qui force le SGBD à faire un 'Full Table Scan' (lecture complète de la table). L'instruction EXPLAIN permet de lire dans les pensées de l'optimiseur de requêtes."
    )

    # Slide 7 (Devenu 8): Algèbre Relationnelle
    add_slide_with_content(
        prs,
        "6. L'Algèbre Relationnelle",
        [
            {"title": "Opérations sur 1 table", "bullets": [
                "Sélection : Filtre les lignes selon une condition.",
                "Projection : Ne conserve que certaines colonnes.",
                "Renommage : Modifie le nom d'une colonne."
            ]},
            {"title": "Opérations sur 2 tables", "highlight": True, "bullets": [
                "Produit Cartésien : Combine tous les éléments de deux tables.",
                "Jointure : Combine des éléments répondant à un critère de liaison."
            ]}
        ],
        "[VISUEL KADEA]\nInfographie montrant le concept visuel de la Projection (colonnes verticales en surbrillance) et de la Sélection (lignes horizontales en surbrillance)."
    )

    # Slide 8: Modélisation UML
    add_slide_with_content(
        prs,
        "7. Modélisation Conceptuelle avec UML",
        [
            {"title": "Diagramme de Classes UML", "bullets": [
                "Le nom de la relation en haut.",
                "La liste des colonnes en bas : nom_colonne: type(longueur).",
                "Les Clés Primaires sont soulignées, les Clés Étrangères précédées d'un #."
            ]},
            {"title": "Agrégation et Composition", "highlight": True, "bullets": [
                "Agrégation : Lien simple (ex: une maison est une agrégation de murs).",
                "Composition : Agrégation forte, la destruction du tout détruit la partie (ex: école et cycles scolaires)."
            ]}
        ],
        "[VISUEL KADEA]\nExemple de formalisme UML strict avec une table d'association.",
        speaker_notes="Rappelez que l'UML est universel. Si un architecte produit un bon diagramme de classes UML, n'importe quel développeur ou outil de génération de code peut le traduire automatiquement en script SQL DDL. La composition (losange plein) implique une dépendance existentielle forte : si je supprime une Commande, ses Lignes_de_Commande doivent être détruites en cascade (ON DELETE CASCADE)."
    )

    # Slide 9 (Devenu 10): SQL DDL
    add_slide_with_content(
        prs,
        "8. Le Langage SQL : Définition des Données",
        [
            {"title": "Architecture Client-Serveur", "bullets": ["Le client (ex: DBeaver) envoie des requêtes SQL au serveur (ex: PostgreSQL)."]},
            {"title": "Création de Tables (DDL)", "highlight": True, "bullets": [
                "Syntaxe : CREATE TABLE nom (colonne1 type, colonne2 type);",
                "Contraintes : NOT NULL, PRIMARY KEY, UNIQUE, CHECK(condition)."
            ]},
            {"title": "Modification", "bullets": ["ALTER TABLE : Ajouter, supprimer des colonnes ou des contraintes.", "DROP TABLE : Supprimer une table."]}
        ],
        "[VISUEL KADEA]\nCapture d'écran de l'interface DBeaver connectée à une base PostgreSQL."
    )

    # Slide 10: SQL DML
    add_slide_with_content(
        prs,
        "9. Le Langage SQL : Manipulation des Données",
        [
            {"title": "Insertion", "bullets": ["INSERT INTO table (col1, col2) VALUES (val1, val2);"]},
            {"title": "Modification", "bullets": ["UPDATE table SET col1 = val1 WHERE condition;"]},
            {"title": "Suppression", "bullets": ["DELETE FROM table WHERE condition;"]},
            {"title": "Qualité de Données (LinkedIn Learning)", "highlight": True, "bullets": [
                "Le SQL est indispensable pour le nettoyage (Data Wrangling).",
                "Permet de gérer massivement les formats mixtes ou les données manquantes (NULL)."
            ]}
        ],
        "[VISUEL KADEA]\nCode snippet : Une requête UPDATE corrigeant une faute de frappe dans un dataset sale."
    )

    # Slide 11 (Devenu 12): SQL DQL (Sélection)
    add_slide_with_content(
        prs,
        "10. Le Langage SQL : Requêtes d'Analyse (1/2)",
        [
            {"title": "Projection et Sélection", "bullets": [
                "SELECT col1, col2 FROM table WHERE condition;",
                "Utilisation de l'opérateur LIKE avec '%' pour les recherches textuelles."
            ]},
            {"title": "La Jointure (JOIN)", "highlight": True, "bullets": [
                "Pour relier les données de plusieurs tables.",
                "SELECT * FROM tableA JOIN tableB ON tableA.id = tableB.fk_id;"
            ]}
        ],
        "[VISUEL KADEA]\nUn schéma d'ensembles de Venn montrant le fonctionnement d'une Jointure (l'intersection)."
    )

    # Slide 12 (Devenu 13): SQL DQL (Agrégation)
    add_slide_with_content(
        prs,
        "12. L'Analyse Data avec SQL : Agrégations",
        [
            {"title": "Agrégations Classiques", "bullets": [
                "COUNT(), AVG(), SUM(), MIN(), MAX().",
                "Le GROUP BY : Réduction de dimensionnalité pour les KPIs.",
                "Le HAVING : Filtrer *après* l'agrégation (contrairement au WHERE)."
            ]}
        ],
        "[VISUEL KADEA]\nConcept du GROUP BY: Des données brutes de différentes couleurs séparées dans des 'seaux' (Buckets) colorés.",
        speaker_notes="Ne confondez jamais WHERE et HAVING. Le WHERE filtre les lignes avant le calcul. Le HAVING filtre le résultat du calcul (ex: afficher les catégories dont la moyenne est > 100)."
    )

    # NOUVELLE SLIDE : SQL AVANCÉ (CTEs et Window Functions)
    add_slide_with_content(
        prs,
        "13. Advanced SQL : CTEs & Window Functions",
        [
            {"title": "Common Table Expressions (CTEs)", "bullets": [
                "La clause WITH : Remplace les sous-requêtes imbriquées illisibles.",
                "Rend le code SQL modulaire, lisible et maintenable."
            ]},
            {"title": "Window Functions (Analyse Statistique)", "highlight": True, "bullets": [
                "Syntaxe : Fonction OVER (PARTITION BY ... ORDER BY ...).",
                "Calculs complexes (Rangs, Totaux cumulés, Moyennes mobiles) *sans* écraser le niveau de détail des lignes.",
                "Fonctions analytiques clés : ROW_NUMBER(), RANK(), LAG(), LEAD()."
            ]}
        ],
        "[VISUEL KADEA]\nComparatif visuel : GROUP BY (écrase les lignes) vs OVER (maintient les lignes d'origine mais ajoute une colonne de calcul).",
        speaker_notes="C'est ici que l'on sépare les juniors des seniors en SQL. Les CTEs (clause WITH) permettent d'écrire des scripts d'analyse clairs comme du code impératif. Les Window Functions sont essentielles en Data Science et BI pour faire du tracking temporel (LAG/LEAD) sans faire de jointures croisées complexes."
    )

    prs.save('Cours_Complet_BDD_et_SQL.pptx')
    print("Présentation complète en français générée avec succès.")

if __name__ == "__main__":
    create_full_presentation()
