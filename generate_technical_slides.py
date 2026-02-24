from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Configuration ---
FILENAME = "Atelier_Technique_Data.pptx"
TITLE_FONT = "Arial"
BODY_FONT = "Calibri"
COLOR_BLUE = RGBColor(0, 51, 153) # Deep Corporate Blue
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ORANGE = RGBColor(255, 140, 0) # Accent Color
COLOR_GRAY = RGBColor(80, 80, 80)
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240)

# --- Helper Functions ---

def create_slide(prs, layout_index=1):
    try:
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        return slide
    except IndexError:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        return slide

def add_title(slide, text, font_size=32, color=COLOR_BLUE):
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.name = TITLE_FONT
        title.text_frame.paragraphs[0].font.size = Pt(font_size)
        title.text_frame.paragraphs[0].font.color.rgb = color
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_content_text(slide, points, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, point in enumerate(points):
        p = tf.add_paragraph()
        p.text = point
        p.font.name = BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_GRAY
        p.space_after = Pt(10)
        p.level = 0

def draw_shape(slide, shape_type, left, top, width, height, color=COLOR_BLUE, text=""):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLOR_GRAY
    if text:
        shape.text = text
        shape.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_code_block(slide, code_text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_GRAY

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = "Consolas"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.LEFT

def add_placeholder_image(slide, description, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 200, 200) # Gray placeholder
    shape.line.dash_style = 4 # Dashed line

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"[Capture d'écran : {description}]"
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

# --- Main Script ---

def generate_presentation():
    prs = Presentation()

    # 1. Slide de Titre
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Atelier Technique : Data Wrangling & Bases de Données"
    title.text_frame.paragraphs[0].font.name = TITLE_FONT
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE

    subtitle.text = "De l'Excel avancé aux opérations CRUD SQL/NoSQL"
    subtitle.text_frame.paragraphs[0].font.name = BODY_FONT
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_ORANGE

    # 2. Agenda
    slide = create_slide(prs, 5)
    add_title(slide, "Au Programme")
    add_content_text(slide, [
        "1. Nettoyage de Données (Excel & Power Query)",
        "2. Tri vs Filtre : Ne plus confondre",
        "3. SQL vs NoSQL : CRUD en pratique",
        "4. Gestion de Données : Bonnes Pratiques",
        "5. Atelier Pratique : Telco & Banque"
    ])

    # 3. Garbage In, Garbage Out
    slide = create_slide(prs, 5)
    add_title(slide, "Le Concept Clé : GIGO")

    draw_shape(slide, MSO_SHAPE.CAN, Inches(1), Inches(3), Inches(1.5), Inches(2), COLOR_GRAY, "Données Sales\n(Garbage In)")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(3), Inches(3.5), Inches(1), Inches(0.5), COLOR_BLUE)
    draw_shape(slide, MSO_SHAPE.GEAR_6, Inches(4.5), Inches(3), Inches(2), Inches(2), COLOR_ORANGE, "Traitement")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(7), Inches(3.5), Inches(1), Inches(0.5), COLOR_BLUE)
    draw_shape(slide, MSO_SHAPE.CAN, Inches(8.5), Inches(3), Inches(1.5), Inches(2), COLOR_GRAY, "Résultat Faux\n(Garbage Out)")

    add_content_text(slide, ["Exemple Telco : Un numéro '+33 6...' vs '06...' crée deux clients différents !"], Inches(1), Inches(5.5), Inches(8))

    # 4. Excel Avancé : Tri vs Filtre
    slide = create_slide(prs, 5)
    add_title(slide, "Excel : Trier ou Filtrer ?")

    draw_shape(slide, MSO_SHAPE.UP_ARROW_CALLOUT, Inches(0.5), Inches(2), Inches(4), Inches(2), COLOR_BLUE, "TRIER (Sort)\nOrganise l'ordre (A-Z, 1-10).\nAucune donnée n'est cachée.\nUsage : Classement, Top 10.")
    draw_shape(slide, MSO_SHAPE.FUNNEL, Inches(5), Inches(2), Inches(4), Inches(2), COLOR_ORANGE, "FILTRER (Filter)\nMasque les lignes non désirées.\nAttention : Les données cachées existent toujours !\nUsage : Focus sur une catégorie.")

    # 5. Power Query (Démo Conceptuelle)
    slide = create_slide(prs, 5)
    add_title(slide, "Au-delà d'Excel : Power Query")
    add_content_text(slide, [
        "Pourquoi ? Excel plante après 1M lignes. Power Query automatise.",
        "Fonctionnalité clé : ETL (Extract, Transform, Load).",
        "Exemple : Séparer 'Nom Prénom' en 2 colonnes."
    ], width=Inches(8))

    add_placeholder_image(slide, "Power Query Editor : Colonne 'Split by Delimiter'", Inches(1), Inches(4), Inches(8), Inches(3))

    # 6. SQL : Structure & CRUD
    slide = create_slide(prs, 5)
    add_title(slide, "SQL : Le Langage Structuré")

    # Table visual
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(3), Inches(2), COLOR_BLUE, "Table 'Clients'\nID | Nom | Solde")

    # CRUD commands
    add_code_block(slide, "CREATE TABLE Clients (id INT, nom VARCHAR(50));\n\n-- C: Create\nINSERT INTO Clients VALUES (1, 'Jean');\n\n-- R: Read\nSELECT * FROM Clients WHERE id=1;\n\n-- U: Update\nUPDATE Clients SET nom='Paul' WHERE id=1;\n\n-- D: Delete\nDELETE FROM Clients WHERE id=1;", Inches(4), Inches(2), Inches(5.5), Inches(5))

    # 7. NoSQL : Flexibilité & CRUD
    slide = create_slide(prs, 5)
    add_title(slide, "NoSQL : L'Approche Document")

    # JSON visual
    add_code_block(slide, "{\n  '_id': 1,\n  'nom': 'Jean',\n  'historique': ['Achat1', 'Achat2']\n}", Inches(0.5), Inches(2), Inches(3.5), Inches(3))

    # CRUD commands
    add_code_block(slide, "// C: Create\ndb.clients.insertOne({nom: 'Jean'});\n\n// R: Read\ndb.clients.find({nom: 'Jean'});\n\n// U: Update\ndb.clients.updateOne({nom: 'Jean'}, {$set: {nom: 'Paul'}});\n\n// D: Delete\ndb.clients.deleteOne({nom: 'Paul'});", Inches(4.5), Inches(2), Inches(5), Inches(5))

    # 8. Concepts Avancés : Delete vs Drop
    slide = create_slide(prs, 5)
    add_title(slide, "Ne pas confondre !")

    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3.5), Inches(2), COLOR_ORANGE, "DELETE\nSupprime les DONNÉES.\nLa table (structure) reste vide.\nOn peut annuler (Rollback).")
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2), Inches(3.5), Inches(2), COLOR_BLUE, "DROP\nSupprime TOUT.\nDonnées + Structure.\nIrréversible (sauf backup).")

    # 9. Standards de Données (ISO)
    slide = create_slide(prs, 5)
    add_title(slide, "Standards Internationaux")

    add_content_text(slide, [
        "Dates (ISO 8601) : YYYY-MM-DD (2023-12-31).",
        "Décimales : Point (.) ou Virgule (,) ? Dépend du système (US vs FR).",
        "Encodage : UTF-8 (pour les accents é, à, ç)."
    ])

    draw_shape(slide, MSO_SHAPE.PLAQUE, Inches(6), Inches(4), Inches(3), Inches(1.5), COLOR_BLUE, "2023-12-31\n✅ Standard")

    # 10. Bonnes Pratiques Gestionnaire
    slide = create_slide(prs, 5)
    add_title(slide, "Checklist du Gestionnaire de Données")

    add_content_text(slide, [
        "1. Toujours garder une copie du fichier BRUT (Raw Data).",
        "2. Ne jamais travailler sur l'original.",
        "3. Documenter les nettoyages faits (ex: 'J'ai remplacé les vides par 0').",
        "4. Vérifier les types (Est-ce que '123' est un nombre ou du texte ?)."
    ])

    prs.save(FILENAME)
    print(f"Presentation saved as {FILENAME}")

if __name__ == "__main__":
    generate_presentation()
