from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Configuration ---
FILENAME = "Atelier_S5_BI.pptx"
TITLE_FONT = "Arial"
BODY_FONT = "Calibri"
COLOR_BLUE = RGBColor(0, 51, 153) # Deep Corporate Blue
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ORANGE = RGBColor(255, 140, 0) # Accent Color
COLOR_GRAY = RGBColor(80, 80, 80)
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240)

# --- Helper Functions ---

def create_slide(prs, layout_index=1):
    try:
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        return slide
    except IndexError:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        return slide

def add_title(slide, text, font_size=32, color=COLOR_BLUE):
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.name = TITLE_FONT
        title.text_frame.paragraphs[0].font.size = Pt(font_size)
        title.text_frame.paragraphs[0].font.color.rgb = color
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_content_text(slide, points, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, point in enumerate(points):
        p = tf.add_paragraph()
        p.text = point
        p.font.name = BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_GRAY
        p.space_after = Pt(10)
        p.level = 0

def draw_shape(slide, shape_type, left, top, width, height, color=COLOR_BLUE, text=""):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLOR_GRAY
    if text:
        shape.text = text
        shape.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

# --- Main Script ---

def generate_presentation():
    prs = Presentation()

    # 1. Slide de Titre
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Semaine 5 : Intro Business Intelligence"
    title.text_frame.paragraphs[0].font.name = TITLE_FONT
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE

    subtitle.text = "Du Tableur Classique au Tableau de Bord Interactif"
    subtitle.text_frame.paragraphs[0].font.name = BODY_FONT
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_ORANGE

    # 2. Théorie : Dashboard vs Excel (La métaphore de la voiture)
    slide = create_slide(prs, 5)
    add_title(slide, "1. Pourquoi un Dashboard ?")

    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(3.5), Inches(4), COLOR_GRAY, "EXCEL\n\nC'est la mécanique.\nDes millions de chiffres.\nIl faut chercher l'info.\nComplexe et statique.")
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5), Inches(2), Inches(4), Inches(4), COLOR_ORANGE, "BI (DASHBOARD)\n\nC'est le tableau de bord.\nUn seul coup d'œil.\nJauges, Alertes, Carte.\nInteractif et Dynamique.")

    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4), Inches(4), Inches(5), Inches(4))

    # 3. Pratique : Atelier Power BI (Transition)
    slide = create_slide(prs, 5)
    add_title(slide, "2. Atelier Pratique : Transition")

    # Processus BI
    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(2), Inches(1.5), COLOR_BLUE, "Connecter\n(Excel, Web...)")
    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(3), Inches(2), Inches(1.5), COLOR_ORANGE, "Transformer\n(Nettoyer)")
    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3), Inches(2), Inches(1.5), COLOR_GRAY, "Visualiser\n(Graphiques)")

    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3), Inches(3.75), Inches(4), Inches(3.75))
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6), Inches(3.75), Inches(7), Inches(3.75))

    add_content_text(slide, ["Exercice : Importer 'dataset_S5_bi.xlsx' dans Power BI Desktop ou Looker Studio."], Inches(1), Inches(5), Inches(8))

    prs.save(FILENAME)
    print(f"Presentation saved as {FILENAME}")

if __name__ == "__main__":
    generate_presentation()
