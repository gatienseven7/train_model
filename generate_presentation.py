from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Couleurs Kadea Academy
KADEA_RED = RGBColor(237, 28, 36)    # #ED1C24
ANTHRACITE = RGBColor(40, 40, 40)    # Noir Anthracite
WHITE = RGBColor(255, 255, 255)      # Blanc

def set_shape_text(shape, text, color=ANTHRACITE, font_size=Pt(18), bold=False, alignment=PP_ALIGN.LEFT):
    """Utility function to set text properties on a shape."""
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.size = font_size
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0] # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    set_shape_text(title, "Module 2: Modélisation Avancée et Solutions Décisionnelles", KADEA_RED, Pt(32), True, PP_ALIGN.CENTER)
    set_shape_text(subtitle, "WEEK 1 | Relational Database Foundations & Schema Design", ANTHRACITE, Pt(24), False, PP_ALIGN.CENTER)

def add_bloc1_slide(prs):
    slide_layout = prs.slide_layouts[5] # Blank layout (for custom layout)
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, "🔴 BLOC 1: CORE CONCEPTS - Le choc des Titans", KADEA_RED, Pt(28), True)

    # Dot-Dash Content (Situation, Complication, Resolution)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p1 = tf.add_paragraph()
    p1.text = "Situation: L'artisanat sur tableur"
    p1.font.bold = True
    p1.font.size = Pt(20)
    p1.font.color.rgb = ANTHRACITE

    p2 = tf.add_paragraph()
    p2.text = "- Excel est excellent, mais possède un plafond de verre : le Excel Ceiling."
    p2.font.size = Pt(16)
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "Complication: Limites de stockage"
    p3.font.bold = True
    p3.font.size = Pt(20)
    p3.font.color.rgb = ANTHRACITE

    p4 = tf.add_paragraph()
    p4.text = "- À plus d'1 million de lignes, le logiciel suffoque.\n- Risque élevé d'erreurs manuelles et de corruption."
    p4.font.size = Pt(16)
    p4.level = 1

    p5 = tf.add_paragraph()
    p5.text = "Résolution: Ingénierie de la donnée avec RDBMS"
    p5.font.bold = True
    p5.font.size = Pt(20)
    p5.font.color.rgb = KADEA_RED

    p6 = tf.add_paragraph()
    p6.text = "- Relational Database Management Systems.\n- Data Integrity: Règles strictes (ex: pas de doublons, atomicité).\n- Scalability: Des milliards de records sans perte de performance."
    p6.font.size = Pt(16)
    p6.level = 1

    # Visual Placeholder
    viz_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = WHITE
    viz_box.line.color.rgb = KADEA_RED
    viz_box.line.width = Pt(2)
    set_shape_text(viz_box, "[VISUEL KADEA]\nBalance minimaliste.\nGauche: Tableur fissuré (Spreadsheet).\nDroite: Cylindre rouge (RDBMS) avec bouclier (Data Integrity).", ANTHRACITE, Pt(14), False, PP_ALIGN.CENTER)

def add_bloc2_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, "🔴 BLOC 2: TECHNICAL VOCAB - Parler le langage machine", KADEA_RED, Pt(28), True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p1 = tf.add_paragraph()
    p1.text = "Tables, Fields & Records"
    p1.font.bold = True
    p1.font.size = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "- Tables: Stockage intuitif des données.\n- Records: Une ligne = une occurrence unique du monde réel.\n- Fields: Colonnes décrivant les attributs de la donnée."
    p2.font.size = Pt(16)
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "Primary Keys (PK) & Foreign Keys (FK)"
    p3.font.bold = True
    p3.font.size = Pt(20)
    p3.font.color.rgb = KADEA_RED

    p4 = tf.add_paragraph()
    p4.text = "- PK: Identifiant absolu et unique de chaque Record.\n- FK: Clé étrangère pointant vers la PK d'une table parente.\n- C'est le ciment de toute base de données mondiale !"
    p4.font.size = Pt(16)
    p4.level = 1

    # Visual Placeholder
    viz_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = WHITE
    viz_box.line.color.rgb = KADEA_RED
    viz_box.line.width = Pt(2)
    set_shape_text(viz_box, "[VISUEL KADEA]\nDiagramme épuré.\nBloc 1: Champ ID rouge vif (PK).\nBloc 2: Champ ID contourné rouge (FK).\nUn rayon laser rouge relie la PK à la FK.", ANTHRACITE, Pt(14), False, PP_ALIGN.CENTER)

def add_bloc3_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, "🔴 BLOC 3: WORKSHOP - Les mains dans le cambouis", KADEA_RED, Pt(28), True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p1 = tf.add_paragraph()
    p1.text = "Setup de l'environnement"
    p1.font.bold = True
    p1.font.size = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "- Installation locale (DBeaver ou SQLite Browser)."
    p2.font.size = Pt(16)
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "Initialisation: Kadea_Lab.db"
    p3.font.bold = True
    p3.font.size = Pt(20)

    p4 = tf.add_paragraph()
    p4.text = "- Table Customers avec PK (ID client).\n- Table Orders avec FK référençant l'ID client."
    p4.font.size = Pt(16)
    p4.level = 1

    p5 = tf.add_paragraph()
    p5.text = "Crash Test: Data Integrity en action"
    p5.font.bold = True
    p5.font.size = Pt(20)
    p5.font.color.rgb = KADEA_RED

    p6 = tf.add_paragraph()
    p6.text = "- Forcez l'insertion d'une commande pour un client inexistant.\n- Résultat : Erreur SGBD. La Data Integrity fonctionne !"
    p6.font.size = Pt(16)
    p6.level = 1

    # Visual Placeholder
    viz_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = WHITE
    viz_box.line.color.rgb = KADEA_RED
    viz_box.line.width = Pt(2)
    set_shape_text(viz_box, "[VISUEL KADEA]\nCapture DBeaver (thème sombre).\nPastilles rouges (1, 2, 3) pour guider le Setup Step-by-Step.", ANTHRACITE, Pt(14), False, PP_ALIGN.CENTER)

def add_bloc4_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, "🔴 BLOC 4: EXERCISE - Design Thinking", KADEA_RED, Pt(28), True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p1 = tf.add_paragraph()
    p1.text = "Conceptual Data Model (CDM/MCD)"
    p1.font.bold = True
    p1.font.size = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "- On ne code jamais avant de modéliser sur papier.\n- Abstraire le 'monde réel' graphiquement."
    p2.font.size = Pt(16)
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "Mission : Kadea Telco (Infrastructure RDC)"
    p3.font.bold = True
    p3.font.size = Pt(20)
    p3.font.color.rgb = KADEA_RED

    p4 = tf.add_paragraph()
    p4.text = "- Modéliser le système d'infrastructures télécoms dans les 26 provinces de la RDC.\n- Entités: Provinces, Cell_Towers (Antennes), Technicians."
    p4.font.size = Pt(16)
    p4.level = 1

    p5 = tf.add_paragraph()
    p5.text = "Le Piège Formateur (Many-to-Many)"
    p5.font.bold = True
    p5.font.size = Pt(20)

    p6 = tf.add_paragraph()
    p6.text = "- Un technicien gère plusieurs antennes, une antenne est gérée par plusieurs techniciens.\n- Solution: Créer une table intermédiaire (ex: Maintenance_Assignments) !"
    p6.font.size = Pt(16)
    p6.level = 1

    # Visual Placeholder
    viz_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = WHITE
    viz_box.line.color.rgb = KADEA_RED
    viz_box.line.width = Pt(2)
    set_shape_text(viz_box, "[VISUEL KADEA]\nAnimation en deux temps.\n1: Croquis crayonné (Provinces/Antennes/Techniciens).\n2: Morphing en Database Schema strict encadré de rouge.", ANTHRACITE, Pt(14), False, PP_ALIGN.CENTER)

def add_acid_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, "🔴 DEEP DIVE: Le Modèle ACID (Transactions)", KADEA_RED, Pt(28), True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p1 = tf.add_paragraph()
    p1.text = "Garantir la fiabilité absolue"
    p1.font.bold = True
    p1.font.size = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "- Atomicité : Une transaction passe entièrement ou est annulée (0 ou 100%).\n- Cohérence : Le Database Schema et les contraintes sont toujours respectés."
    p2.font.size = Pt(16)
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "Sécurité & Isolation"
    p3.font.bold = True
    p3.font.size = Pt(20)
    p3.font.color.rgb = KADEA_RED

    p4 = tf.add_paragraph()
    p4.text = "- Isolation : Les transactions simultanées n'interfèrent pas (Multi-users).\n- Durabilité : En cas de crash, les données validées sont sauvegardées."
    p4.font.size = Pt(16)
    p4.level = 1

    # Visual Placeholder
    viz_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = WHITE
    viz_box.line.color.rgb = KADEA_RED
    viz_box.line.width = Pt(2)
    set_shape_text(viz_box, "[VISUEL KADEA]\nSchéma ACID.\n4 piliers supportant un temple grec (RDBMS).\nChaque pilier porte une lettre A, C, I, D avec une icône de bouclier.", ANTHRACITE, Pt(14), False, PP_ALIGN.CENTER)

def add_sql_analysis_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, "🔴 DEEP DIVE: SQL Syntax pour l'Analyse (DBeaver)", KADEA_RED, Pt(28), True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p1 = tf.add_paragraph()
    p1.text = "Nettoyage & Data Quality"
    p1.font.bold = True
    p1.font.size = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "- Gérer les 'Dirty Data' (doublons, formats mixtes, Missing Values).\n- SQL permet de transformer et filtrer massivement (Data Wrangling)."
    p2.font.size = Pt(16)
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "Analyse Statistique & Window Functions"
    p3.font.bold = True
    p3.font.size = Pt(20)
    p3.font.color.rgb = KADEA_RED

    p4 = tf.add_paragraph()
    p4.text = "- Syntaxe avancée (PostgreSQL) : SELECT, JOIN, GROUP BY.\n- Fonctions de fenêtrage (OVER, PARTITION BY) pour des analyses de distribution."
    p4.font.size = Pt(16)
    p4.level = 1

    # Visual Placeholder
    viz_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = WHITE
    viz_box.line.color.rgb = KADEA_RED
    viz_box.line.width = Pt(2)
    set_shape_text(viz_box, "[VISUEL KADEA]\nSplit Screen.\nHaut : Un dataset brouillon ('Dirty').\nBas : Bloc de code SQL (Syntaxe PostgreSQL/DBeaver) purifiant les données (Entonnoir Rouge).", ANTHRACITE, Pt(14), False, PP_ALIGN.CENTER)

def add_conclusion_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    set_shape_text(title_box, "Ressources & Next Steps", KADEA_RED, Pt(28), True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    p1 = tf.add_paragraph()
    p1.text = "Key Milestone Validé ✅"
    p1.font.bold = True
    p1.font.size = Pt(22)

    p2 = tf.add_paragraph()
    p2.text = "- Limites du tableur comprises.\n- Jargon de l'ingénieur de données acquis.\n- Environnement technique opérationnel (DBeaver/SQLite).\n- Capacité à lire un plan d'architecture (Schema Design)."
    p2.font.size = Pt(18)
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "\nRessources Utiles :"
    p3.font.bold = True
    p3.font.size = Pt(22)
    p3.font.color.rgb = KADEA_RED

    p4 = tf.add_paragraph()
    p4.text = "- Documentation SQLite officielle: sqlite.org/docs.html\n- DBeaver Community Edition: dbeaver.io\n- Visualisation de Database Schema: dbdiagram.io"
    p4.font.size = Pt(18)
    p4.level = 1

def create_presentation():
    prs = Presentation()

    add_title_slide(prs)
    add_bloc1_slide(prs)
    add_bloc2_slide(prs)
    add_bloc3_slide(prs)
    add_bloc4_slide(prs)
    add_acid_slide(prs)
    add_sql_analysis_slide(prs)
    add_conclusion_slide(prs)

    prs.save('Module2_Week1_Presentation.pptx')
    print("Présentation 'Module2_Week1_Presentation.pptx' générée avec succès.")

if __name__ == "__main__":
    create_presentation()
