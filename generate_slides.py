from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Configuration ---
FILENAME = "The_Data_Journey.pptx"
TITLE_FONT = "Arial"
BODY_FONT = "Calibri"
COLOR_RED = RGBColor(220, 20, 60) # Crimson Red
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_GRAY = RGBColor(105, 105, 105)

# --- Helper Functions ---

def create_slide(prs, layout_index=1):
    """Creates a new slide based on layout index."""
    try:
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        return slide
    except IndexError:
        # Fallback if layout index is out of range
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        return slide


def add_title(slide, text, font_size=36, color=COLOR_RED):
    """Adds or modifies the title of a slide."""
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.name = TITLE_FONT
        title.text_frame.paragraphs[0].font.size = Pt(font_size)
        title.text_frame.paragraphs[0].font.color.rgb = color
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_content_text(slide, points, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(5)):
    """Adds a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        p = tf.add_paragraph()
        p.text = point
        p.font.name = BODY_FONT
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_BLACK
        p.space_after = Pt(14)
        if i == 0 and len(points) == 1: # Single line
            p.level = 0
        else:
            p.level = 0

def draw_shape(slide, shape_type, left, top, width, height, color=COLOR_RED, text=""):
    """Draws a geometric shape."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLOR_GRAY

    if text:
        shape.text = text
        shape.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

# --- Main Script ---

def generate_presentation():
    prs = Presentation()

    # 1. Slide de Titre
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide Layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "The Data Journey – De l'Excel au Big Data"
    title.text_frame.paragraphs[0].font.name = TITLE_FONT
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_RED

    subtitle.text = "Du tableur à l'architecture de données moderne."
    subtitle.text_frame.paragraphs[0].font.name = BODY_FONT
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_GRAY

    # Visuel Slide 1: Rectangle -> Cylindre
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.5), Inches(1.5), Inches(1), COLOR_RED, "Excel")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(4.7), Inches(4.8), Inches(1), Inches(0.4), COLOR_GRAY)
    draw_shape(slide, MSO_SHAPE.CAN, Inches(6), Inches(4.2), Inches(1.5), Inches(1.5), COLOR_RED, "Base de Données")


    # 2. L'Ancrage
    slide = create_slide(prs, 5) # Blank layout
    add_title(slide, "Où en sommes-nous ?")
    add_content_text(slide, [
        "Semaine 1-2 : Saisie de données.",
        "Semaine 3 : Organisation & Tri (Excel).",
        "Objectif : Devenir Gestionnaire de Données."
    ], width=Inches(5))

    # Visuel Slide 2: Bonhomme et Tableau
    draw_shape(slide, MSO_SHAPE.SMILEY_FACE, Inches(6.5), Inches(2), Inches(1.5), Inches(1.5), COLOR_RED)
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(4), Inches(0.5), Inches(0.5), COLOR_GRAY)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(5), Inches(2.5), Inches(1.5), COLOR_RED, "Data")


    # 3. La Question Fatidique (Comparatif)
    slide = create_slide(prs, 5)
    add_title(slide, "Est-ce qu'Excel est une Base de Données ?")

    # Colonne Excel
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(2.5), Inches(4), COLOR_WHITE).line.color.rgb = COLOR_RED
    add_content_text(slide, ["EXCEL", "Oui, < 1M lignes", "Analyse Visuelle", "Risque: Pas de sécurité"], Inches(0.6), Inches(2.1), Inches(2.3), Inches(3.5))

    # Colonne SQL
    draw_shape(slide, MSO_SHAPE.CAN, Inches(3.5), Inches(2), Inches(2.5), Inches(4), COLOR_RED)
    txBox = slide.shapes.add_textbox(Inches(3.6), Inches(2.5), Inches(2.3), Inches(3))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "SQL (Relationnel)\n\nLa Rigueur\nSécurité\nCohérence Bancaire"
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Colonne NoSQL
    draw_shape(slide, MSO_SHAPE.CLOUD, Inches(6.5), Inches(2), Inches(3), Inches(4), COLOR_GRAY)
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(3), Inches(2.4), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "NoSQL (Big Data)\n\nLa Vitesse\nVolume Massif\nFlexibilité"
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER


    # 4. L'Évolution de la Data
    slide = create_slide(prs, 5)
    add_title(slide, "De l'Archive Papier au Cloud")

    # Timeline
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(0.5), Inches(4), Inches(9), Inches(0.5), COLOR_GRAY)

    # Jalons
    draw_shape(slide, MSO_SHAPE.FOLDED_CORNER, Inches(1), Inches(2.5), Inches(1.5), Inches(1.2), COLOR_RED, "Avant 1960\nPapier")
    draw_shape(slide, MSO_SHAPE.CAN, Inches(4), Inches(2.5), Inches(1.5), Inches(1.2), COLOR_RED, "1970\nSQL (Banque)")
    draw_shape(slide, MSO_SHAPE.CLOUD, Inches(7), Inches(2.5), Inches(1.5), Inches(1.2), COLOR_RED, "2000+\nNoSQL (Web)")


    # 5. SQL
    slide = create_slide(prs, 5)
    add_title(slide, "Le Monde Carré du SQL")
    add_content_text(slide, [
        "Structure : Tables, Lignes, Colonnes.",
        "Force : Intégrité (ACID).",
        "Usage : Finance, RH, ERP."
    ], width=Inches(5))

    # Visuel Tables reliées
    t1 = draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(2), Inches(1.5), Inches(2), COLOR_RED, "Table A")
    t2 = draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8), Inches(4), Inches(1.5), Inches(2), COLOR_GRAY, "Table B")
    # Using specific connector type
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.75), Inches(3), Inches(8.75), Inches(5))


    # 6. NoSQL - Documents
    slide = create_slide(prs, 5)
    add_title(slide, "NoSQL : Orienté Documents")
    add_content_text(slide, [
        "Format : JSON / XML (Hiérarchique).",
        "Usage : Profils utilisateurs, Web.",
        "Exemple : MongoDB."
    ], width=Inches(5))
    draw_shape(slide, MSO_SHAPE.FOLDED_CORNER, Inches(6.5), Inches(2.5), Inches(2), Inches(2.5), COLOR_RED, "{ JSON }")

    # 7. NoSQL - Clé-Valeur
    slide = create_slide(prs, 5)
    add_title(slide, "NoSQL : Clé-Valeur")
    add_content_text(slide, [
        "Concept : Une étiquette = Une donnée.",
        "Usage : Paniers, Sessions, Cache.",
        "Exemple : Redis."
    ], width=Inches(5))
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(3), Inches(1), Inches(1), COLOR_GRAY, "KEY")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(7.1), Inches(3.2), Inches(0.8), Inches(0.5), COLOR_BLACK)
    draw_shape(slide, MSO_SHAPE.CUBE, Inches(8), Inches(2.5), Inches(1.5), Inches(1.5), COLOR_RED, "VALUE")


    # 8. NoSQL - Colonnes
    slide = create_slide(prs, 5)
    add_title(slide, "NoSQL : Orienté Colonnes")
    add_content_text(slide, [
        "Concept : Lecture massive d'une info.",
        "Usage : Logs Telco, Bourse.",
        "Exemple : Cassandra."
    ], width=Inches(5))
    # Cylindre tranché
    draw_shape(slide, MSO_SHAPE.CAN, Inches(7), Inches(2), Inches(1.5), Inches(4), COLOR_RED)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(2.5), Inches(1.3), Inches(0.1), COLOR_WHITE)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(3.5), Inches(1.3), Inches(0.1), COLOR_WHITE)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(4.5), Inches(1.3), Inches(0.1), COLOR_WHITE)


    # 9. NoSQL - Graphes
    slide = create_slide(prs, 5)
    add_title(slide, "NoSQL : Orienté Graphes")
    add_content_text(slide, [
        "Concept : Tout est connecté.",
        "Usage : Réseaux sociaux, Fraude.",
        "Exemple : Neo4j."
    ], width=Inches(5))
    # Reseau
    c1 = draw_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(2), Inches(0.8), Inches(0.8), COLOR_RED)
    c2 = draw_shape(slide, MSO_SHAPE.OVAL, Inches(8), Inches(3), Inches(0.8), Inches(0.8), COLOR_RED)
    c3 = draw_shape(slide, MSO_SHAPE.OVAL, Inches(7), Inches(5), Inches(0.8), Inches(0.8), COLOR_RED)
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.4), Inches(2.4), Inches(8.4), Inches(3.4))
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.4), Inches(3.4), Inches(7.4), Inches(5.4))
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(7.4), Inches(5.4), Inches(6.4), Inches(2.4))


    # 10. Pipeline Data
    slide = create_slide(prs, 5)
    add_title(slide, "Le Voyage de la Donnée")

    draw_shape(slide, MSO_SHAPE.CHEVRON, Inches(0.5), Inches(3), Inches(2), Inches(1), COLOR_GRAY, "1. Collecte")
    draw_shape(slide, MSO_SHAPE.CHEVRON, Inches(2.8), Inches(3), Inches(2), Inches(1), COLOR_GRAY, "2. Stockage")
    draw_shape(slide, MSO_SHAPE.CHEVRON, Inches(5.1), Inches(3), Inches(2), Inches(1), COLOR_RED, "3. ETL")
    draw_shape(slide, MSO_SHAPE.CHEVRON, Inches(7.4), Inches(3), Inches(2), Inches(1), COLOR_GRAY, "4. Visu")


    # 11. CRUD
    slide = create_slide(prs, 5)
    add_title(slide, "Les 4 Commandements (CRUD)")

    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(1.5), Inches(1.5), COLOR_RED, "C\nCreate")
    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(3), Inches(1.5), Inches(1.5), COLOR_RED, "R\nRead")
    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(3), Inches(1.5), Inches(1.5), COLOR_RED, "U\nUpdate")
    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3), Inches(1.5), Inches(1.5), COLOR_RED, "D\nDelete")


    # 12. Standards
    slide = create_slide(prs, 5)
    add_title(slide, "Parler le Langage Machine")
    add_content_text(slide, [
        "ISO 8601 : YYYY-MM-DD.",
        "Atomicité : Une cellule = Une info.",
        "Pas de 'Paris, France' -> Deux colonnes !"
    ], width=Inches(6))
    draw_shape(slide, MSO_SHAPE.PLAQUE, Inches(7), Inches(3), Inches(2), Inches(1.5), COLOR_RED, "YYYY-MM-DD\n✅")


    # 13. Identifiant Unique
    slide = create_slide(prs, 5)
    add_title(slide, "Le Piège des Noms")
    add_content_text(slide, [
        "Problème : Homonymes (Jean Dupont).",
        "Solution : Clé Primaire (ID #12345).",
        "Règle : Ne jamais se fier au texte."
    ], width=Inches(5))
    draw_shape(slide, MSO_SHAPE.SMILEY_FACE, Inches(6), Inches(3), Inches(1), Inches(1), COLOR_GRAY)
    draw_shape(slide, MSO_SHAPE.SMILEY_FACE, Inches(7.5), Inches(3), Inches(1), Inches(1), COLOR_GRAY)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(4.2), Inches(1), Inches(0.5), COLOR_RED, "#001")
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(4.2), Inches(1), Inches(0.5), COLOR_RED, "#002")


    # 14. Trier vs Filtrer
    slide = create_slide(prs, 5)
    add_title(slide, "Ne Confondez Pas !")

    # Trier
    draw_shape(slide, MSO_SHAPE.UP_ARROW_CALLOUT, Inches(1), Inches(3), Inches(3), Inches(2), COLOR_GRAY, "TRIER\nOrganiser l'ordre.\nNe cache rien.")
    # Filtrer
    draw_shape(slide, MSO_SHAPE.FUNNEL, Inches(5), Inches(3), Inches(2), Inches(2), COLOR_RED, "FILTRER\nMasquer.\nRisque d'oubli !")


    # 15. Garbage In Garbage Out
    slide = create_slide(prs, 5)
    add_title(slide, "Garbage In, Garbage Out")

    draw_shape(slide, MSO_SHAPE.CAN, Inches(1), Inches(3), Inches(1.5), Inches(2), COLOR_GRAY, "Trash")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(3), Inches(3.5), Inches(1), Inches(0.5), COLOR_RED)
    draw_shape(slide, MSO_SHAPE.GEAR_6, Inches(4.5), Inches(3), Inches(2), Inches(2), COLOR_BLACK, "Process")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(7), Inches(3.5), Inches(1), Inches(0.5), COLOR_RED)
    draw_shape(slide, MSO_SHAPE.CAN, Inches(8.5), Inches(3), Inches(1.5), Inches(2), COLOR_GRAY, "Trash")


    # 16. Danger Zone
    slide = create_slide(prs, 5)
    add_title(slide, "Danger Zone")
    add_content_text(slide, [
        "Cellules fusionnées (Cauchemar).",
        "Formats mixtes (12 vs '12').",
        "Pas de Backup."
    ], width=Inches(6))
    draw_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(7), Inches(2), Inches(2), Inches(2), COLOR_RED, "!")


    # 17. Ressources
    slide = create_slide(prs, 5)
    add_title(slide, "Pour Aller Plus Loin")
    add_content_text(slide, [
        "Doc Microsoft Excel.",
        "Tutos SQL pour débutants.",
        "Bonnes pratiques de nettoyage."
    ])


    # 18. Questions
    slide = create_slide(prs, 5)
    add_title(slide, "Questions ?")
    draw_shape(slide, MSO_SHAPE.BALLOON, Inches(4), Inches(3), Inches(3), Inches(2), COLOR_RED, "?")


    prs.save(FILENAME)
    print(f"Presentation saved as {FILENAME}")

if __name__ == "__main__":
    generate_presentation()
