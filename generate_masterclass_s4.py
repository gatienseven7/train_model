from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Configuration ---
FILENAME = "Masterclass_S4_Logique_Viz.pptx"
TITLE_FONT = "Arial"
BODY_FONT = "Calibri"
COLOR_BLUE_DARK = RGBColor(15, 32, 67) # Very dark professional blue
COLOR_ACCENT = RGBColor(211, 47, 47) # Sharp Red for emphasis
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT = RGBColor(50, 50, 50)
COLOR_LIGHT_GRAY = RGBColor(245, 245, 245)

# --- Helper Functions ---

def create_slide(prs, layout_index=5): # Default to blank layout
    try:
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        return slide
    except IndexError:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        return slide

def add_title(slide, text, subtitle="", font_size=36, color=COLOR_BLUE_DARK):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = TITLE_FONT
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = BODY_FONT
        p2.font.size = Pt(20)
        p2.font.color.rgb = COLOR_ACCENT
        p2.font.italic = True

def add_content(slide, title, points, left=Inches(0.5), top=Inches(2), width=Inches(8.5), height=Inches(5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = TITLE_FONT
        p_title.font.size = Pt(24)
        p_title.font.color.rgb = COLOR_BLUE_DARK
        p_title.font.bold = True
        p_title.space_after = Pt(14)

    for i, point in enumerate(points):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.text = point
        p.font.name = BODY_FONT
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)
        p.level = 0 if title else 1 # Indent if there is a sub-title above

def draw_shape(slide, shape_type, left, top, width, height, fill_color, text="", text_color=COLOR_WHITE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color # Clean border
    if text:
        shape.text = text
        shape.text_frame.paragraphs[0].font.name = BODY_FONT
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        shape.text_frame.paragraphs[0].font.color.rgb = text_color
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_pro_tip(slide, text, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(9), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    shape.line.color.rgb = COLOR_ACCENT
    shape.line.width = Pt(2)

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "PRO-TIP: " + text
    p.font.name = BODY_FONT
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.LEFT

# --- Main Script ---

def generate_presentation():
    prs = Presentation()

    # --- Slide 1: Hook & Vision ---
    slide = create_slide(prs, 5) # Blank
    add_title(slide, "Au-delà des chiffres : L'intelligence de la donnée", "Passer de la 'Data Brute' à l''Information Actionnable'")

    add_content(slide, "", [
        "La donnée n'a de valeur que si elle déclenche une action.",
        "Nous passons du statut de 'Saisisseur' à celui de 'Décideur'.",
        "La Data Brute est muette : Un tableau de 10 000 lignes ne raconte rien sans contexte.",
        "L'Intelligence commence ici : La logique conditionnelle automatise la pensée.",
        "La Visualisation est le traducteur : Elle convertit la complexité en évidence stratégique."
    ], top=Inches(2), width=Inches(9))

    # Flow visual
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.5), Inches(1.8), Inches(1), COLOR_LIGHT_GRAY, "Donnée\n(Brute)", COLOR_TEXT)
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(3), Inches(4.8), Inches(0.5), Inches(0.4), COLOR_BLUE_DARK)
    draw_shape(slide, MSO_SHAPE.GEAR_6, Inches(3.7), Inches(4.5), Inches(1), Inches(1), COLOR_BLUE_DARK, "Logique\n(SI)")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(4.9), Inches(4.8), Inches(0.5), Inches(0.4), COLOR_BLUE_DARK)
    draw_shape(slide, MSO_SHAPE.CHART_X, Inches(5.6), Inches(4.5), Inches(1.2), Inches(1), COLOR_BLUE_DARK, "Viz")
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(7), Inches(4.8), Inches(0.5), Inches(0.4), COLOR_ACCENT)
    draw_shape(slide, MSO_SHAPE.OVAL, Inches(7.7), Inches(4.4), Inches(1.2), Inches(1.2), COLOR_ACCENT, "Décision\n(Action)")


    # --- Slide 2: Logique Conditionnelle ---
    slide = create_slide(prs, 5)
    add_title(slide, "Le premier algorithme : Automatiser la décision avec '=SI()'")

    add_content(slide, "L'Anatomie d'une règle", [
        "Syntaxe : =SI(Test logique ; Action si Vrai ; Action si Faux)"
    ], top=Inches(1.8))

    add_content(slide, "Impact Métier (Use Cases)", [
        "RH : Automatisation du statut 'Admis/Recalé' sans relecture humaine.",
        "Finance : Alerte automatique en cas de 'Dépassement de Budget'.",
        "Commerce : Déclenchement d'un bonus si 'Objectif Atteint'."
    ], top=Inches(3))

    # Visual Code Block
    draw_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(2), Inches(3.5), Inches(1.5), COLOR_BLUE_DARK, "IF Budget > 10k\nTHEN Alert = 'ROUGE'\nELSE Alert = 'OK'")

    add_pro_tip(slide, "Une bonne condition est binaire : C'est Oui ou Non. C'est l'essence même de la logique informatique et du reporting par exception.", Inches(5.5))


    # --- Slide 3: Choix Graphique ---
    slide = create_slide(prs, 5)
    add_title(slide, "L'Art du Choix Graphique : La rigueur contre le superflu")

    add_content(slide, "", ["Un mauvais graphique ment. Le choix du type de visualisation dicte la vérité transmise."], top=Inches(1.5))

    # 3 Columns Compare
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.8), Inches(2), COLOR_LIGHT_GRAY, "HISTOGRAMME (Barres)\n\nBenchmark.\nComparer des catégories (Vendeur A vs B).", COLOR_TEXT)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(3.6), Inches(2.5), Inches(2.8), Inches(2), COLOR_LIGHT_GRAY, "COURBE (Lignes)\n\nTendance.\nSérie temporelle, évolution (Saisonnalité).", COLOR_TEXT)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(2.5), Inches(2.8), Inches(2), COLOR_LIGHT_GRAY, "CAMEMBERT (Secteurs)\n\nProportion.\nMix de coûts (Max 2 à 5 catégories).", COLOR_TEXT)

    add_pro_tip(slide, "La Règle d'Edward Tufte : Le Data-Ink Ratio. Effacez tout ce qui n'est pas de la donnée (grilles inutiles, effets 3D). Maximisez le message.", Inches(5.5))


    # --- Slide 4: Anatomie d'un Graphique ---
    slide = create_slide(prs, 5)
    add_title(slide, "Déconstruire pour convaincre : Du X au Y")

    add_content(slide, "Les Axes ne sont pas juste mathématiques", [
        "Axe Abscisses (X) = Le Contexte : Segmentation (Vendeurs), Temps (Mois). Le 'Qui' ou le 'Quand'.",
        "Axe Ordonnées (Y) = La Valeur Métier : Performance (CA), Taux de conversion. Le 'Combien'.",
        "Storytelling : Ne montrez pas 'Le CA par Vendeur'. Racontez 'Pourquoi l'équipe Sud surperforme au T3'."
    ], top=Inches(1.5), width=Inches(5))

    # Minimalist Chart Visual
    # Y Axis
    draw_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(6), Inches(4.5), Inches(2), Inches(0.1), COLOR_BLUE_DARK) # X Axis line
    draw_shape(slide, MSO_SHAPE.UP_ARROW, Inches(6), Inches(2.5), Inches(0.1), Inches(2), COLOR_BLUE_DARK) # Y Axis line
    # Bars
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(3.5), Inches(0.4), Inches(1), COLOR_BLUE_DARK)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.8), Inches(0.4), Inches(1.7), COLOR_ACCENT)
    draw_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(3.8), Inches(0.4), Inches(0.7), COLOR_BLUE_DARK)

    # Annotations
    tx1 = slide.shapes.add_textbox(Inches(5), Inches(2.5), Inches(1), Inches(0.5))
    tx1.text_frame.text = "Axe Y\n(Impact)"
    tx1.text_frame.paragraphs[0].font.size = Pt(12)
    tx1.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    tx2 = slide.shapes.add_textbox(Inches(7), Inches(4.6), Inches(1.5), Inches(0.5))
    tx2.text_frame.text = "Axe X (Contexte)"
    tx2.text_frame.paragraphs[0].font.size = Pt(12)
    tx2.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT


    # --- Slide 5: Checklist Expert ---
    slide = create_slide(prs, 5)
    add_title(slide, "La Checklist du Data Storyteller", "La visualisation est une discipline de simplification extrême.")

    add_content(slide, "", [
        "1. Le Titre est une conclusion : Pas 'Ventes 2023', mais 'Croissance record portée par le B2B en 2023'.",
        "2. Data-Ink Ratio absolu : Supprimez les quadrillages denses, ombres portées et couleurs agressives inutiles.",
        "3. Étiquetage explicite : Des axes toujours nommés, des unités claires (K€, %). Ne forcez pas l'audience à deviner.",
        "4. Zéro Biais Visuel : Fuyez la 3D. L'axe Y d'un histogramme doit toujours commencer à Zéro pour ne pas exagérer les écarts."
    ], top=Inches(2), width=Inches(9))

    # Quote
    quote = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf_q = quote.text_frame
    p_q = tf_q.paragraphs[0]
    p_q.text = "\"Le but de la visualisation est l'insight, pas l'image.\" - Ben Shneiderman"
    p_q.font.name = TITLE_FONT
    p_q.font.size = Pt(22)
    p_q.font.italic = True
    p_q.font.color.rgb = COLOR_BLUE_DARK
    p_q.alignment = PP_ALIGN.CENTER

    prs.save(FILENAME)
    print(f"Presentation saved as {FILENAME}")

if __name__ == "__main__":
    generate_presentation()
