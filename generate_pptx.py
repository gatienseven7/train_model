from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()

    # Define color scheme (Red and White corporate theme)
    KADEA_RED = RGBColor(220, 38, 38)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(51, 51, 51)

    # 1. TITLE SLIDE (SCR Framework - Situation/Complication/Resolution)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # Background shape for branding
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = KADEA_RED
    bg_shape.line.color.rgb = KADEA_RED

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "PROJET FIL ROUGE : KADEA TELCO"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.bold = True

    subtitle.text = ("L'enfer des requêtes et de la rétention\n\n"
                     "[Situation] : Vous êtes Data Analyst face à une concurrence féroce.\n"
                     "[Complication] : L'instabilité réseau menace la survie de l'entreprise.\n"
                     "[Résolution] : Manipuler et croiser des milliers de données brutes "
                     "pour piloter l'activité avec des règles métier complexes.")

    # Method to format slide titles
    def add_branded_title(slide, text):
        title = slide.shapes.title
        title.text = text
        title.text_frame.paragraphs[0].font.color.rgb = KADEA_RED
        title.text_frame.paragraphs[0].font.bold = True

        # Add a decorative red line under the title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(8), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = KADEA_RED
        line.line.fill.background()

    # Phase Slides using 'Dot-Dash' (Key message title supported by concrete bullet points)
    bullet_slide_layout = prs.slide_layouts[1]

    # PHASE 1
    slide1 = prs.slides.add_slide(bullet_slide_layout)
    add_branded_title(slide1, "PHASE 1 : Le chaos des Logs")
    tf1 = slide1.placeholders[1].text_frame
    tf1.text = "Relier des tables hétérogènes sans erreur de correspondance"

    p1_1 = tf1.add_paragraph()
    p1_1.text = "Scénario : Réconcilier les logs bruts avec le registre de maintenance."
    p1_1.level = 1

    p1_2 = tf1.add_paragraph()
    p1_2.text = "Compétences : Manipulation de matrices, gestion des erreurs (#N/A)."
    p1_2.level = 1

    p1_3 = tf1.add_paragraph()
    p1_3.text = "Outils : RECHERCHEV / VLOOKUP, Nettoyage de données."
    p1_3.level = 1

    p1_4 = tf1.add_paragraph()
    p1_4.text = "Livrable : Une base unifiée traduisant les codes obscurs en adresses."
    p1_4.level = 1

    # PHASE 2
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    add_branded_title(slide2, "PHASE 2 : La synthèse macroscopique")
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Passer de la donnée unitaire à la statistique agrégée"

    p2_1 = tf2.add_paragraph()
    p2_1.text = "Scénario : Offrir au management une vision hélicoptère pour envoyer les réparateurs."
    p2_1.level = 1

    p2_2 = tf2.add_paragraph()
    p2_2.text = "Compétences : Fonctions d'agrégation, analyse de zones critiques."
    p2_2.level = 1

    p2_3 = tf2.add_paragraph()
    p2_3.text = "Outils : GROUP BY, Tableaux Croisés Dynamiques (TCD), SUM, AVG."
    p2_3.level = 1

    p2_4 = tf2.add_paragraph()
    p2_4.text = "Livrable : Tableau de synthèse multicritères (Région / Technologie)."
    p2_4.level = 1

    # PHASE 3
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    add_branded_title(slide3, "PHASE 3 : Le labyrinthe de la Rétention")
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Traduire une règle métier complexe en formule mathématique"

    p3_1 = tf3.add_paragraph()
    p3_1.text = "Scénario : Coder un plan de rétention complexe pour limiter le churn des clients."
    p3_1.level = 1

    p3_2 = tf3.add_paragraph()
    p3_2.text = "Compétences : Logique algorithmique avancée, calcul d'impact financier."
    p3_2.level = 1

    p3_3 = tf3.add_paragraph()
    p3_3.text = "Outils : Fonctions SI (IF) imbriquées, croisées avec opérateurs ET (AND) / OU (OR)."
    p3_3.level = 1

    p3_4 = tf3.add_paragraph()
    p3_4.text = "Livrable : Colonne 'Taux de Remise' 100% automatisée."
    p3_4.level = 1

    # PHASE 4
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    add_branded_title(slide4, "PHASE 4 : Industrialisation via BI")
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Passer du fichier statique au Dashboard dynamique"

    p4_1 = tf4.add_paragraph()
    p4_1.text = "Scénario : Abandonner le tableur lourd pour un pilotage visuel, fluide et en temps réel."
    p4_1.level = 1

    p4_2 = tf4.add_paragraph()
    p4_2.text = "Compétences : Processus ETL, modélisation de données, Design interactif."
    p4_2.level = 1

    p4_3 = tf4.add_paragraph()
    p4_3.text = "Outils : Power BI (Migration, relations natives, segments dynamiques)."
    p4_3.level = 1

    p4_4 = tf4.add_paragraph()
    p4_4.text = "Livrable : Tableau de bord complet avec segments et cartes géographiques."
    p4_4.level = 1

    prs.save('Presentation_Projet.pptx')
    print("Presentation_Projet.pptx generated.")

if __name__ == '__main__':
    create_presentation()
